import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import re

class ExcelToPPTGenerator:
    def __init__(self, base_path):
        self.base_path = base_path
        self.input_json_path = os.path.join(base_path, "input", "csv", "input.json")
        self.preview_json_path = os.path.join(base_path, "output", "preview.json")
        self.font_alignment_path = os.path.join(base_path, "output", "font_alignment.json")
        self.templates_path = os.path.join(base_path, "input", "templates")
        self.cleaned_templates_path = os.path.join(base_path, "input", "cleaned_templates")
        self.default_templates_path = os.path.join(base_path, "default_templates")
        
        # Create necessary directories
        os.makedirs(self.cleaned_templates_path, exist_ok=True)
        os.makedirs(os.path.join(base_path, "output"), exist_ok=True)
    
    def load_json(self, filepath):
        """Load JSON file"""
        with open(filepath, 'r', encoding='utf-8') as f:
            return json.load(f)
    
    def hex_to_rgb(self, hex_color):
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def clean_template(self, template_path, cleaned_path):
        """Remove only text content, keep all design elements"""
        prs = Presentation(template_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    try:
                        shape.text_frame.clear()
                    except Exception as e:
                        print(f"Could not clear text from shape: {e}")
                
                if shape.is_placeholder:
                    try:
                        if hasattr(shape, "text_frame"):
                            shape.text_frame.clear()
                    except Exception as e:
                        print(f"Could not clear placeholder text: {e}")
        
        prs.save(cleaned_path)
        print(f"✅ Template cleaned - text removed, design preserved")
        return len(prs.slides)
    
    def get_font_size_for_slide(self, slide_index, font_alignment_data, default_title=28, default_content=16):
        """Get fixed font sizes from alignment data"""
        if not font_alignment_data:
            return default_title, default_content, None
        
        font_mapping = font_alignment_data.get("font_size_mapping", [])
        
        for slide_info in font_mapping:
            if slide_info.get("slide_index") == slide_index:
                title_font = slide_info.get("title_font_size", default_title)
                content_font = slide_info.get("content_font_size", default_content)
                subtitle_font = slide_info.get("subtitle_font_size")
                return title_font, content_font, subtitle_font
        
        return default_title, default_content, None
    
    def add_title_to_slide(self, slide, title, font_color_rgb, slide_width, slide_height, 
                          title_font_size, is_first_slide=False, is_last_slide=False):
        """Add title with FIXED font size"""
        if is_first_slide:
            left = Inches(0.5)
            top = slide_height / 2 - Inches(1.2)
            width = slide_width - Inches(1)
            height = Inches(1.5)
        elif is_last_slide:
            left = Inches(0.5)
            top = slide_height / 2 - Inches(1.0)
            width = slide_width - Inches(1)
            height = Inches(1.2)
        else:
            left = Inches(0.5)
            top = Inches(0.4)
            width = slide_width - Inches(1)
            height = Inches(0.8)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = title
        
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.line_spacing = 1.2
        
        run = paragraph.runs[0]
        run.font.size = Pt(title_font_size)  # FIXED SIZE from LLM
        run.font.bold = True
        run.font.color.rgb = RGBColor(*font_color_rgb)
        
        text_frame.word_wrap = True
        if is_first_slide or is_last_slide:
            text_frame.vertical_anchor = 1
    
    def add_subtitle_to_slide(self, slide, subtitle, font_color_rgb, slide_width, slide_height, 
                             subtitle_font_size, is_first_slide=False):
        """Add subtitle with FIXED font size"""
        if is_first_slide:
            left = slide_width / 2 + Inches(0.5)
            top = slide_height / 2 + Inches(0.3)
            width = slide_width / 2 - Inches(1)
            height = Inches(0.6)
            alignment = PP_ALIGN.LEFT
        else:
            left = Inches(0.5)
            top = slide_height / 2 + Inches(0.2)
            width = slide_width - Inches(1)
            height = Inches(0.6)
            alignment = PP_ALIGN.CENTER
        
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = subtitle_box.text_frame
        text_frame.text = subtitle
        
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = alignment
        
        run = paragraph.runs[0]
        run.font.size = Pt(subtitle_font_size)  # FIXED SIZE from LLM
        run.font.italic = True
        run.font.color.rgb = RGBColor(*font_color_rgb)
        
        text_frame.word_wrap = True
        if not is_first_slide:
            text_frame.vertical_anchor = 1
    
    def add_content_to_slide(self, slide, content, font_color_rgb, left, top, width, height, 
                            content_font_size, has_image=False):
        """Add content with FIXED font size - NO dynamic calculation"""
        content_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = content_box.text_frame
        
        # Clean content - remove any markdown symbols
        if isinstance(content, list):
            content_text = '\n\n'.join(str(item).strip() for item in content)
        else:
            content_text = str(content).strip()
        
        # Remove ** markdown symbols
        content_text = content_text.replace('**', '')
        
        text_frame.text = content_text
        
        # Apply FIXED font size from LLM
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.line_spacing = 1.3
            paragraph.space_before = Pt(6)
            paragraph.space_after = Pt(3)
            
            for run in paragraph.runs:
                run.font.size = Pt(content_font_size)  # FIXED SIZE
                run.font.color.rgb = RGBColor(*font_color_rgb)
        
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        
        print(f"   Content: {len(content_text)} chars, Font: {content_font_size}pt (FIXED)")
        
        return content_font_size
    
    def add_image_to_slide(self, slide, image_path, left, top, width, height, is_dashboard=False):
        """Add image with proper aspect ratio"""
        if not os.path.exists(image_path):
            print(f"   ⚠️ Image not found: {image_path}")
            return False
        
        try:
            img = Image.open(image_path)
            img_width, img_height = img.size
            aspect_ratio = img_width / img_height
            
            target_width = width
            target_height = target_width / aspect_ratio
            
            if target_height > height:
                target_height = height
                target_width = target_height * aspect_ratio
            
            # Apply ratio constraint
            if is_dashboard:
                desired_ratio = 4/3
            else:
                desired_ratio = 2/3
            
            current_ratio = target_width / target_height
            
            if current_ratio > desired_ratio:
                target_height = target_width / desired_ratio
                if target_height > height:
                    target_height = height
                    target_width = target_height * desired_ratio
            else:
                target_width = target_height * desired_ratio
                if target_width > width:
                    target_width = width
                    target_height = target_width / desired_ratio
            
            left_offset = left + (width - target_width) / 2
            top_offset = top + (height - target_height) / 2
            
            slide.shapes.add_picture(image_path, left_offset, top_offset, target_width, target_height)
            print(f"   ✅ Image added: {target_width/Inches(1):.1f}\" x {target_height/Inches(1):.1f}\"")
            return True
            
        except Exception as e:
            print(f"   ❌ Error adding image: {e}")
            return False
    
    def get_template_path(self, template_name):
        """Get template path - custom first, then default"""
        custom_template_path = os.path.join(self.templates_path, template_name)
        if os.path.exists(custom_template_path):
            print(f"Using custom template: {template_name}")
            return custom_template_path
        
        default_templates = ['Creative', 'Professional', 'Minimal', 'Technical']
        template_lower = template_name.lower()
        
        for default_template in default_templates:
            if template_lower == default_template.lower():
                default_template_path = os.path.join(self.default_templates_path, f"{default_template}.pptx")
                if os.path.exists(default_template_path):
                    print(f"Using default template: {default_template}")
                    return default_template_path
        
        fallback_path = os.path.join(self.default_templates_path, "Professional.pptx")
        if os.path.exists(fallback_path):
            print(f"Using fallback template: Professional")
            return fallback_path
        
        raise FileNotFoundError(f"No template found for {template_name}")
    
    def generate_ppt(self):
        """Main function to generate PowerPoint with FIXED font sizes"""
        
        # Load input files
        input_data = self.load_json(self.input_json_path)
        preview_data = self.load_json(self.preview_json_path)
        
        # Load font alignment data if available
        font_alignment_data = None
        if os.path.exists(self.font_alignment_path):
            font_alignment_data = self.load_json(self.font_alignment_path)
            print("✅ Font alignment data loaded - using FIXED font sizes")
        else:
            print("⚠️ No font alignment data - using default sizes")
        
        # Get template info
        presentation_title = input_data['presentation_title']
        template_name = input_data['template_name']
        font_color_hex = input_data['text_color']
        font_color_rgb = self.hex_to_rgb(font_color_hex)
        
        # Get template path and clean it
        template_path = self.get_template_path(template_name)
        cleaned_template_path = os.path.join(self.cleaned_templates_path, os.path.basename(template_path))
        
        print(f"\n🎨 Cleaning template: {template_name}")
        num_template_slides = self.clean_template(template_path, cleaned_template_path)
        
        # Load cleaned template
        prs = Presentation(cleaned_template_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        print(f"📐 Slide dimensions: {slide_width/Inches(1):.1f}\" x {slide_height/Inches(1):.1f}\"")
        
        # Calculate slides needed
        total_slides_needed = len(preview_data['slides'])
        print(f"📊 Slides needed: {total_slides_needed}, Template has: {num_template_slides}")
        
        # Add slides if needed
        if num_template_slides < total_slides_needed:
            slides_to_add = total_slides_needed - num_template_slides
            print(f"➕ Adding {slides_to_add} slides...")
            
            original_slides = list(prs.slides)
            
            for i in range(slides_to_add):
                ppt_slide_index = num_template_slides + i
                
                if ppt_slide_index == 0:
                    template_slide_index = 0
                elif ppt_slide_index == total_slides_needed - 1:
                    template_slide_index = num_template_slides - 1
                else:
                    if num_template_slides > 2:
                        template_slide_index = (ppt_slide_index - 1) % (num_template_slides - 2) + 1
                    else:
                        template_slide_index = min(ppt_slide_index, num_template_slides - 1)
                
                if template_slide_index < len(original_slides):
                    template_slide = original_slides[template_slide_index]
                    new_slide = prs.slides.add_slide(template_slide.slide_layout)
                    
                    for shape in template_slide.shapes:
                        if hasattr(shape, 'element'):
                            new_slide.shapes._spTree.append(shape.element)
        
        print(f"\n📝 Generating {total_slides_needed} slides with FIXED fonts...\n")
        
        # Process each slide
        for slide_data in preview_data['slides']:
            slide_index = slide_data['slide_index']
            print(f"🔹 Slide {slide_index}:")
            
            if slide_index > len(prs.slides):
                print(f"   ⚠️ Slide {slide_index} exceeds available slides")
                continue
            
            slide = prs.slides[slide_index - 1]
            placeholders = slide_data['placeholders']
            
            # Get FIXED font sizes from LLM alignment data
            title_font_size, content_font_size, subtitle_font_size = self.get_font_size_for_slide(
                slide_index, font_alignment_data
            )
            
            # Determine slide type
            is_first_slide = slide_index == 1
            is_last_slide = slide_index == total_slides_needed
            
            # Check dashboard slide
            title_text = placeholders.get('title', '').lower() if not is_first_slide else ''
            is_dashboard_slide = 'dashboard' in title_text
            
            # Calculate body area
            body_top = Inches(1.5)
            body_height = slide_height - Inches(2.0)
            body_width = slide_width - Inches(1.2)
            
            if is_first_slide or is_last_slide:
                body_top = slide_height / 2 + Inches(0.5)
                body_height = slide_height / 2 - Inches(1.0)
            
            # Add title
            if is_first_slide:
                title_text = presentation_title
            else:
                title_text = placeholders.get('title', '')
            
            if title_text:
                self.add_title_to_slide(slide, title_text, font_color_rgb, slide_width, slide_height,
                                       title_font_size, is_first_slide, is_last_slide)
                print(f"   Title: \"{title_text[:50]}...\" Font: {title_font_size}pt")
            
            # Handle subtitles
            if 'subtitle' in placeholders:
                subtitle = placeholders['subtitle']
                if not subtitle_font_size:
                    subtitle_font_size = 18 if is_first_slide else 20
                
                self.add_subtitle_to_slide(slide, subtitle, font_color_rgb, slide_width, slide_height,
                                          subtitle_font_size, is_first_slide)
                print(f"   Subtitle: Font: {subtitle_font_size}pt")
                
                if is_first_slide or is_last_slide:
                    continue
            
            # Get content
            content_text = placeholders.get('content', '')
            image_path = placeholders.get('image_path', '')
            
            if not content_text and not image_path:
                continue
            
            # Dashboard slide
            if is_dashboard_slide and image_path and os.path.exists(image_path):
                dashboard_left = Inches(0.5)
                dashboard_top = Inches(1.2)
                dashboard_width = slide_width - Inches(1)
                dashboard_height = slide_height - Inches(2.0)
                
                self.add_image_to_slide(slide, image_path, dashboard_left, dashboard_top, 
                                       dashboard_width, dashboard_height, True)
                continue
            
            # Handle content and images
            if image_path and os.path.exists(image_path):
                # Split layout: content + image
                content_width = body_width * 0.5
                image_width = body_width * 0.45
                gap = Inches(0.3)
                
                # Add content on left
                if content_text:
                    content_left = Inches(0.5)
                    self.add_content_to_slide(slide, content_text, font_color_rgb,
                                            content_left, body_top, content_width, body_height,
                                            content_font_size, has_image=True)
                
                # Add image on right
                image_left = Inches(0.5) + content_width + gap
                self.add_image_to_slide(slide, image_path, image_left, body_top, 
                                       image_width, body_height, is_dashboard_slide)
            
            elif content_text:
                # Full width content
                content_left = Inches(0.5)
                self.add_content_to_slide(slide, content_text, font_color_rgb,
                                        content_left, body_top, body_width, body_height,
                                        content_font_size, has_image=False)
        
        # Save presentation
        output_path = os.path.join(self.base_path, "output", "generated_presentation.pptx")
        prs.save(output_path)
        
        print(f"\n{'='*60}")
        print(f"✅ Presentation saved: {output_path}")
        print(f"{'='*60}")
        
        return output_path


# Main execution
if __name__ == "__main__":
    # Set your base path
    base_path = r"C:\Users\anith\Downloads\backend (2)\backend"  # ✅ YOUR PATH
    
    generator = ExcelToPPTGenerator(base_path)
    output_file = generator.generate_ppt()
    print(f"\n🎉 Successfully generated: {output_file}")
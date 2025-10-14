import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import colorsys
import re
import sys

class ExcelToPPTGenerator:
    def __init__(self, base_path=None):
        # Base path defaults to current directory/backend for compatibility with your layout
        if base_path:
            self.base_path = base_path
        else:
            # Try common folder names, otherwise use cwd
            candidate = os.path.join(os.getcwd(), "backend")
            self.base_path = candidate if os.path.exists(candidate) else os.getcwd()

        self.input_json_path = os.path.join(self.base_path, "input", "csv", "input.json")
        # Use preview_final.json produced by service.py
        self.preview_json_path = os.path.join(self.base_path, "output", "preview_final.json")
        self.templates_path = os.path.join(self.base_path, "input", "templates")
        self.cleaned_templates_path = os.path.join(self.base_path, "input", "cleaned_templates")
        self.default_templates_path = os.path.join(self.base_path, "default_templates")

        # Create necessary directories
        os.makedirs(self.cleaned_templates_path, exist_ok=True)
        os.makedirs(os.path.join(self.base_path, "output"), exist_ok=True)

    def load_json(self, filepath):
        """Load JSON file (safe)"""
        if not os.path.exists(filepath):
            raise FileNotFoundError(f"Required JSON not found: {filepath}")
        with open(filepath, 'r', encoding='utf-8') as f:
            return json.load(f)

    def hex_to_rgb(self, hex_color):
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def get_contrasting_color(self, bg_color_rgb):
        """Generate a contrasting color based on background color"""
        r, g, b = [x / 255.0 for x in bg_color_rgb]
        h, s, v = colorsys.rgb_to_hsv(r, g, b)

        if v < 0.5:
            new_v = 0.9
        else:
            new_v = 0.1

        new_r, new_g, new_b = colorsys.hsv_to_rgb(h, s, new_v)
        return (int(new_r * 255), int(new_g * 255), int(new_b * 255))

    def extract_slide_background_color(self, slide):
        """Extract background color from slide"""
        try:
            if slide.background.fill.type == 1:  # Solid fill
                color = slide.background.fill.fore_color.rgb
                return (color[0], color[1], color[2])
        except Exception:
            pass
        # Default white background if cannot extract
        return (255, 255, 255)

    def clean_template(self, template_path, cleaned_path):
        """Remove only text content from text boxes and placeholders, keep all design elements"""
        prs = Presentation(template_path)

        for slide in prs.slides:
            for shape in slide.shapes:
                # Only clear text content, don't remove the shape itself
                if hasattr(shape, "text_frame"):
                    try:
                        text_frame = shape.text_frame
                        text_frame.clear()
                    except Exception as e:
                        print(f"Could not clear text from shape: {e}")

                # If it's a placeholder, clear its text but keep the placeholder
                if getattr(shape, "is_placeholder", False):
                    try:
                        if hasattr(shape, "text_frame"):
                            shape.text_frame.clear()
                    except Exception as e:
                        print(f"Could not clear placeholder text: {e}")

        prs.save(cleaned_path)
        print(f"Template cleaned - cleared text content, kept all design elements")
        return len(prs.slides)

    def calculate_dynamic_font_size(self, text, allocated_width_inches, allocated_height_inches, has_image=False):
        """Calculate optimal font size based on text length and allocated space"""
        if isinstance(text, list):
            text_length = sum(len(str(item)) for item in text)
        else:
            text_length = len(str(text))

        area = allocated_width_inches * allocated_height_inches
        if area <= 0:
            return 12

        chars_per_sq_inch = text_length / area

        if text_length < 50:
            font_size = 34
        elif text_length < 100:
            font_size = 32
        elif text_length < 150:
            font_size = 30
        elif text_length < 250:
            font_size = 28
        elif text_length < 350:
            font_size = 26
        elif text_length < 500:
            font_size = 24
        elif text_length < 700:
            font_size = 22
        elif text_length < 900:
            font_size = 20
        elif text_length < 1200:
            font_size = 18
        else:
            font_size = 16

        if allocated_height_inches < 1.0:
            font_size = min(font_size, 12)
        elif allocated_height_inches < 1.5:
            font_size = min(font_size, 14)
        elif allocated_height_inches < 2.0:
            font_size = min(font_size, 16)
        elif allocated_height_inches < 2.5:
            font_size = min(font_size, 18)
        elif allocated_height_inches < 3.0:
            font_size = min(font_size, 20)

        if allocated_width_inches < 2.0:
            font_size = min(font_size, 14)
        elif allocated_width_inches < 3.0:
            font_size = min(font_size, 16)

        if has_image:
            font_size = max(10, font_size - 2)

        font_size = max(10, min(font_size, 36))

        print(f"Text adaptation: {text_length} chars, {area:.1f} sq.in -> font: {font_size}pt")
        return font_size

    def add_title_to_slide(self, slide, title, font_color_rgb, slide_width, slide_height, is_first_slide=False, is_last_slide=False):
        """Add title to slide header with improved alignment"""
        if is_first_slide:
            left = Inches(0.5)
            top = slide_height / 2 - Inches(1.2)
            width = slide_width - Inches(1)
            height = Inches(1.5)
        elif is_last_slide:
            left = Inches(0.5)
            top = slide_height / 2 - Inches(1.0)
            width = slide_width - Inches(1)
            height = Inches(1.2)
        else:
            left = Inches(0.5)
            top = Inches(0.4)
            width = slide_width - Inches(1)
            height = Inches(0.8)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = title

        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.line_spacing = 1.2

        run = paragraph.runs[0]
        if is_first_slide:
            run.font.size = Pt(36)
        elif is_last_slide:
            run.font.size = Pt(36)
        else:
            run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*font_color_rgb)

        text_frame.word_wrap = True
        text_frame.auto_size = True
        if is_first_slide or is_last_slide:
            text_frame.vertical_anchor = 1

    def add_subtitle_to_first_slide(self, slide, subtitle, font_color_rgb, slide_width, slide_height):
        """Add subtitle to first slide at right side from center"""
        left = slide_width / 2 + Inches(0.5)
        top = slide_height / 2 + Inches(0.3)
        width = slide_width / 2 - Inches(1)
        height = Inches(0.6)

        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = subtitle_box.text_frame
        text_frame.text = subtitle

        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT

        run = paragraph.runs[0]
        run.font.size = Pt(18)
        run.font.italic = True
        run.font.color.rgb = RGBColor(*font_color_rgb)

        text_frame.word_wrap = True

    def add_subtitle_to_last_slide(self, slide, subtitle, font_color_rgb, slide_width, slide_height):
        """Add subtitle to last slide centered below title"""
        left = Inches(0.5)
        top = slide_height / 2 + Inches(0.2)
        width = slide_width - Inches(1)
        height = Inches(0.6)

        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = subtitle_box.text_frame
        text_frame.text = subtitle

        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER

        run = paragraph.runs[0]
        run.font.size = Pt(20)
        run.font.italic = True
        run.font.color.rgb = RGBColor(*font_color_rgb)

        text_frame.word_wrap = True
        text_frame.vertical_anchor = 1

    def add_content_to_slide(self, slide, content, font_color_rgb, left, top, width, height, has_image=False):
        """Add content to slide body with improved dynamic font sizing"""
        content_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = content_box.text_frame

        if isinstance(content, list):
            content_text = '\n'.join(str(item) for item in content)
        else:
            content_text = str(content)

        text_frame.text = content_text

        width_inches = width / Inches(1)
        height_inches = height / Inches(1)
        font_size = self.calculate_dynamic_font_size(content_text, width_inches, height_inches, has_image)

        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.line_spacing = 1.2 + (font_size / 100)

            if font_size <= 14:
                paragraph.space_before = Pt(3)
                paragraph.space_after = Pt(1)
            elif font_size <= 18:
                paragraph.space_before = Pt(4)
                paragraph.space_after = Pt(2)
            else:
                paragraph.space_before = Pt(6)
                paragraph.space_after = Pt(3)

            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                run.font.color.rgb = RGBColor(*font_color_rgb)

        text_frame.word_wrap = True

        margin_base = font_size * 0.1
        text_frame.margin_left = Inches(min(0.15, margin_base / 100))
        text_frame.margin_right = Inches(min(0.15, margin_base / 100))
        text_frame.margin_top = Inches(min(0.05, margin_base / 200))
        text_frame.margin_bottom = Inches(min(0.05, margin_base / 200))

        return font_size

    def add_image_to_slide(self, slide, image_path, left, top, width, height, is_dashboard=False):
        """Add image to slide with dynamic sizing in 2:3 ratio or 3:1 for dashboard"""
        if not os.path.exists(image_path):
            print(f"Warning: Image not found - {image_path}")
            return False

        try:
            img = Image.open(image_path)
            img_width, img_height = img.size
            aspect_ratio = img_width / img_height

            target_width = width
            target_height = target_width / aspect_ratio

            if target_height > height:
                target_height = height
                target_width = target_height * aspect_ratio

            if is_dashboard:
                desired_ratio = 4/3
            else:
                desired_ratio = 2/3

            current_ratio = target_width / target_height

            if current_ratio > desired_ratio:
                target_height = target_width / desired_ratio
                if target_height > height:
                    target_height = height
                    target_width = target_height * desired_ratio
            else:
                target_width = target_height * desired_ratio
                if target_width > width:
                    target_width = width
                    target_height = target_width / desired_ratio

            left_offset = left + (width - target_width) / 2
            top_offset = top + (height - target_height) / 2

            slide.shapes.add_picture(image_path, left_offset, top_offset, target_width, target_height)
            ratio_type = "4:3" if is_dashboard else "2:3"
            print(f"Added image with size: {target_width / Inches(1):.2f}\" x {target_height / Inches(1):.2f}\" (ratio: {ratio_type})")
            return True
        except Exception as e:
            print(f"Error adding image {image_path}: {e}")
            return False

    def add_multiple_images_to_slide(self, slide, image_paths, left, top, width, height, is_dashboard=False):
        """Add multiple images in grid format with improved spacing"""
        valid_images = []

        for img_path in image_paths:
            if os.path.exists(img_path):
                valid_images.append(img_path)
            else:
                print(f"Warning: Image not found - {img_path}")

        if not valid_images:
            print("No valid images found to display")
            return

        num_images = len(valid_images)

        if num_images == 1:
            rows, cols = 1, 1
        elif num_images == 2:
            rows, cols = 1, 2
        elif num_images <= 4:
            rows, cols = 2, 2
        elif num_images <= 6:
            rows, cols = 2, 3
        else:
            rows, cols = 3, 3

        cell_width = width / cols
        cell_height = height / rows

        horizontal_padding = Inches(0.15)
        vertical_padding = Inches(0.1)

        img_width = cell_width - horizontal_padding * 2
        img_height = cell_height - vertical_padding * 2

        for idx, img_path in enumerate(valid_images[:rows * cols]):
            row = idx // cols
            col = idx % cols

            img_left = left + col * cell_width + horizontal_padding
            img_top = top + row * cell_height + vertical_padding

            self.add_image_to_slide(slide, img_path, img_left, img_top, img_width, img_height, is_dashboard)

    def extract_image_paths_from_content(self, content_text):
        """Extract image paths from content text"""
        image_paths = []

        if isinstance(content_text, list):
            content_string = ' '.join(str(item) for item in content_text)
        else:
            content_string = str(content_text)

        img_pattern = r'\(Chart:\s*([^)]+)\)'
        matches = re.findall(img_pattern, content_string)

        for match in matches:
            img_relative = match.strip()

            img_relative_normalized = img_relative.replace('\\', os.sep).replace('/', os.sep)

            if not os.path.isabs(img_relative_normalized):
                full_path = os.path.join(self.base_path, img_relative_normalized)
            else:
                full_path = img_relative_normalized

            if os.path.exists(full_path):
                if full_path not in image_paths:
                    image_paths.append(full_path)
            else:
                print(f"Warning: Image not found at path: {full_path}")

        clean_content = re.sub(img_pattern, '', content_string).strip()

        return image_paths, clean_content

    def get_template_path(self, template_name):
        """Get the correct template path - custom templates have priority, then default templates"""
        custom_template_path = os.path.join(self.templates_path, template_name)
        if os.path.exists(custom_template_path):
            print(f"Using custom template: {template_name}")
            return custom_template_path

        default_templates = ['Creative', 'Professional', 'Minimal', 'Technical']
        template_lower = template_name.lower()
        for default_template in default_templates:
            if template_lower == default_template.lower():
                default_template_path = os.path.join(self.default_templates_path, f"{default_template}.pptx")
                if os.path.exists(default_template_path):
                    print(f"Using default template: {default_template}")
                    return default_template_path

        fallback_template = "Professional.pptx"
        fallback_path = os.path.join(self.default_templates_path, fallback_template)
        if os.path.exists(fallback_path):
            print(f"Template {template_name} not found, using fallback template: {fallback_template}")
            return fallback_path
        else:
            raise FileNotFoundError(f"Template {template_name} not found and no fallback template available")

    def generate_ppt(self):
        """Main function to generate PowerPoint presentation"""
        # Load input files
        try:
            input_data = self.load_json(self.input_json_path)
        except FileNotFoundError:
            print(f"Input JSON not found at {self.input_json_path}. Proceeding with minimal defaults.")
            input_data = {
                "presentation_title": "Presentation",
                "template_name": "Professional.pptx",
                "text_color": "#000000"
            }

        try:
            preview_data = self.load_json(self.preview_json_path)
        except FileNotFoundError:
            raise FileNotFoundError(f"preview_final.json not found at {self.preview_json_path}. Run service.py first.")

        presentation_title = input_data.get('presentation_title', 'Presentation')
        template_name = input_data.get('template_name', 'Professional.pptx')
        font_color_hex = input_data.get('text_color', '#000000')
        font_color_rgb = self.hex_to_rgb(font_color_hex)

        template_path = self.get_template_path(template_name)
        cleaned_template_path = os.path.join(self.cleaned_templates_path, os.path.basename(template_path))

        print(f"Cleaning template: {template_name}")
        num_template_slides = self.clean_template(template_path, cleaned_template_path)
        print(f"Template has {num_template_slides} slides")

        prs = Presentation(cleaned_template_path)

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        print(f"Slide dimensions: {slide_width / Inches(1):.2f}\" x {slide_height / Inches(1):.2f}\"")

        if len(prs.slides) > 0:
            bg_color = self.extract_slide_background_color(prs.slides[0])
            print(f"Detected background color: RGB{bg_color}")
            contrasting_color = self.get_contrasting_color(bg_color)
            print(f"Using contrasting color: RGB{contrasting_color}")
        else:
            contrasting_color = font_color_rgb

        total_slides_needed = len(preview_data.get('slides', []))

        print(f"Total slides needed: {total_slides_needed}")
        print(f"Template slides available: {num_template_slides}")

        if num_template_slides >= total_slides_needed:
            print(f"Using existing {num_template_slides} template slides")
        else:
            slides_to_add = total_slides_needed - num_template_slides
            print(f"Adding {slides_to_add} additional slides...")
            original_slides = list(prs.slides)

            for i in range(slides_to_add):
                ppt_slide_index = num_template_slides + i
                if ppt_slide_index == 0:
                    template_slide_index = 0
                elif ppt_slide_index == total_slides_needed - 1:
                    template_slide_index = num_template_slides - 1
                else:
                    if num_template_slides > 2:
                        template_slide_index = (ppt_slide_index - 1) % (num_template_slides - 2) + 1
                    else:
                        template_slide_index = min(ppt_slide_index, num_template_slides - 1)

                if template_slide_index < len(original_slides):
                    template_slide = original_slides[template_slide_index]
                    new_slide = prs.slides.add_slide(template_slide.slide_layout)
                    for shape in template_slide.shapes:
                        if hasattr(shape, 'element'):
                            new_slide.shapes._spTree.append(shape.element)
                    print(f"PPT Slide {ppt_slide_index + 1} -> Template Slide {template_slide_index + 1}")
                else:
                    new_slide = prs.slides.add_slide(prs.slide_layouts[0])
                    print(f"PPT Slide {ppt_slide_index + 1} -> Default Layout")

        print(f"\nGenerating {total_slides_needed} slides with preserved designs...")

        for slide_data in preview_data.get('slides', []):
            slide_index = slide_data.get('slide_index', 1) - 1
            if slide_index >= len(prs.slides):
                print(f"Warning: Slide index {slide_index + 1} exceeds available slides")
                continue

            slide = prs.slides[slide_index]
            placeholders = slide_data.get('placeholders', {})

            is_first_slide = slide_index == 0
            is_last_slide = slide_index == total_slides_needed - 1

            title_text = placeholders.get('title', '') if not is_first_slide else ''
            is_dashboard_slide = 'dashboard' in title_text.lower() if title_text else False

            body_top = Inches(1.5)
            body_height = slide_height - Inches(2.0)
            body_width = slide_width - Inches(1.2)

            if is_first_slide or is_last_slide:
                body_top = slide_height / 2 + Inches(0.5)
                body_height = slide_height / 2 - Inches(1.0)

            if is_first_slide:
                title_text = presentation_title
            else:
                title_text = placeholders.get('title', '')

            if title_text:
                self.add_title_to_slide(slide, title_text, font_color_rgb, slide_width, slide_height, is_first_slide, is_last_slide)

            if is_first_slide and 'subtitle' in placeholders:
                subtitle = placeholders['subtitle']
                self.add_subtitle_to_first_slide(slide, subtitle, font_color_rgb, slide_width, slide_height)
                continue

            if is_last_slide and 'subtitle' in placeholders:
                subtitle = placeholders['subtitle']
                self.add_subtitle_to_last_slide(slide, subtitle, font_color_rgb, slide_width, slide_height)
                continue

            if is_last_slide and 'content' not in placeholders:
                continue

            content_text = placeholders.get('content', '')

            image_path = placeholders.get('image_path', '')

            if is_dashboard_slide and image_path and os.path.exists(image_path):
                dashboard_left = Inches(0.5)
                dashboard_top = Inches(1.2)
                dashboard_width = slide_width - Inches(1)
                dashboard_height = slide_height - Inches(2.0)

                image_added = self.add_image_to_slide(slide, image_path, dashboard_left, dashboard_top, dashboard_width, dashboard_height, True)
                if not image_added:
                    print(f"Slide {slide_index + 1}: Dashboard image not found at {image_path}")
                continue

            if not content_text:
                continue

            if image_path and os.path.exists(image_path):
                image_paths = [image_path]
                clean_content = content_text
            else:
                image_paths, clean_content = self.extract_image_paths_from_content(content_text)

            if isinstance(clean_content, list):
                content_length = sum(len(str(item)) for item in clean_content)
            else:
                content_length = len(str(clean_content))

            if image_paths and clean_content:
                if content_length < 100:
                    content_width = body_width * 0.4
                    image_width = body_width * 0.55
                elif content_length < 300:
                    content_width = body_width * 0.5
                    image_width = body_width * 0.45
                else:
                    content_width = body_width * 0.6
                    image_width = body_width * 0.35

                gap = Inches(0.3)

                content_left = Inches(0.5)
                content_top_adjusted = body_top + Inches(0.1)
                font_size = self.add_content_to_slide(slide, clean_content, font_color_rgb, content_left, content_top_adjusted, content_width, body_height - Inches(0.1), has_image=True)

                image_left = Inches(0.5) + content_width + gap
                image_top_adjusted = body_top + Inches(0.1)
                if len(image_paths) == 1:
                    image_added = self.add_image_to_slide(slide, image_paths[0], image_left, image_top_adjusted, image_width, body_height - Inches(0.1), is_dashboard_slide)
                    if not image_added:
                        print(f"Slide {slide_index + 1}: Image space left empty (image not found)")
                else:
                    self.add_multiple_images_to_slide(slide, image_paths, image_left, image_top_adjusted, image_width, body_height - Inches(0.1), is_dashboard_slide)

                print(f"Slide {slide_index + 1}: Content length {content_length} chars -> font: {font_size}pt")
            elif image_paths and not clean_content:
                image_left = Inches(0.5)
                image_top_adjusted = body_top + Inches(0.1)
                if len(image_paths) == 1:
                    self.add_image_to_slide(slide, image_paths[0], image_left, image_top_adjusted, body_width, body_height - Inches(0.1), is_dashboard_slide)
                else:
                    self.add_multiple_images_to_slide(slide, image_paths, image_left, image_top_adjusted, body_width, body_height - Inches(0.1), is_dashboard_slide)
                print(f"Slide {slide_index + 1}: Images only layout")
            elif clean_content and not image_paths:
                content_left = Inches(0.5)
                content_top_adjusted = body_top + Inches(0.1)
                font_size = self.add_content_to_slide(slide, clean_content, font_color_rgb, content_left, content_top_adjusted, body_width, body_height - Inches(0.1), has_image=False)
                print(f"Slide {slide_index + 1}: Content length {content_length} chars -> font: {font_size}pt (full width)")

        output_path = os.path.join(self.base_path, "output", "generated_presentation.pptx")
        prs.save(output_path)
        print(f"\nPresentation saved to: {output_path}")
        return output_path


if __name__ == "__main__":
    # If you want to override base path, set environment variable EXCEL_PPT_BASE_PATH or supply path as arg
    arg_path = None
    if len(sys.argv) > 1:
        arg_path = sys.argv[1]
    base_path = arg_path or os.environ.get("EXCEL_PPT_BASE_PATH", None)
    generator = ExcelToPPTGenerator(base_path)
    output_file = generator.generate_ppt()
    print(f"Successfully generated: {output_file}")
